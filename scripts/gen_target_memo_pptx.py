#!/usr/bin/env python3
"""
Generate a Target Memo PPTX using the Moat in the Machine template.
Usage: python3 gen_target_memo_pptx.py data.json output.pptx [template.pptx]

v3 -- 10-slide structure (Mar 2026).
Slide order: Title, Snapshot, Napkin, Company Scores, Software Moat,
             Data Moat & Platform, Competitive, GenAI, 100-Day Playbook, Sources.
Handles optional JSON keys: playbook_100day_structured, partnership_dependency,
pricing_model_risk, v2_delta.
"""
import json, sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Resolve skill directory for bundled template ──
SKILL_DIR = Path(__file__).resolve().parent.parent
DEFAULT_TEMPLATE = SKILL_DIR / "templates" / "moat_in_the_machine_v2.pptx"

# ── Template design system (from moat_in_the_machine_v2.pptx) ──
FONT       = "Open Sauce One"
DARK_NAVY  = RGBColor(0x0F, 0x1B, 0x2D)
ACCENT_RED = RGBColor(0xE8, 0x56, 0x3A)
GRAY       = RGBColor(0x5A, 0x65, 0x77)
LIGHT_GRAY = RGBColor(0x8B, 0x95, 0xA5)
FOOTER_CLR = RGBColor(0xFF, 0xFF, 0xFF)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREEN      = RGBColor(0x2E, 0x7D, 0x32)
BLUE       = RGBColor(0x15, 0x65, 0xC0)
RED        = RGBColor(0xC6, 0x28, 0x28)
AMBER      = RGBColor(0xF9, 0xA8, 0x25)
CARD_BG    = RGBColor(0xF5, 0xF6, 0xF8)
CARD_BORDER= RGBColor(0xE0, 0xE3, 0xE8)
SUBTITLE_C = RGBColor(0xC5, 0xCC, 0xD6)

QUAD_COLORS = {
    "AI Goldmine": GREEN, "Contrarian Bet": BLUE,
    "Sand Castle": RED,   "Rising Tide": AMBER,
}
QUAD_LABELS = {
    "AI Goldmine": "BUY", "Contrarian Bet": "HOLD / BUILD",
    "Sand Castle": "AVOID", "Rising Tide": "SELECTIVE",
}

# Standard positions (EMU)
LEFT_MARGIN  = 640080   # ~0.7"
CONTENT_W    = 7863840  # ~8.6"
FOOTER_Y     = 4828032
FOOTER_H     = 274320
HALF_W       = 3749040  # half content width
RIGHT_COL    = 4754880  # right column start


def safe(text):
    if not isinstance(text, str): return str(text)
    return (text
        .replace("\u2013", "-").replace("\u2014", "--")
        .replace("\u2018", "'").replace("\u2019", "'")
        .replace("\u201c", '"').replace("\u201d", '"')
        .replace("\u2022", "-").replace("\u2026", "...")
        .replace("\u2265", ">=").replace("\u2264", "<=")
        .replace("\u00d7", "x").replace("\u2248", "~")
        .replace("\u0394", "Delta ").replace("\u03c9", "omega "))


def _add_footer(slide, company_name="", date=""):
    txBox = slide.shapes.add_textbox(457200, FOOTER_Y, 7772400, FOOTER_H)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = safe(f"Target Memo: {company_name}  |  Dutchess Management  |  {date}  |  Confidential")
    p.font.size = Pt(8)
    p.font.name = FONT
    p.font.color.rgb = FOOTER_CLR


def _text(slide, left, top, width, height, text, size=12, bold=False,
          color=DARK_NAVY, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = safe(text)
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT
    p.alignment = align
    return txBox


def _section_header(slide, text, top=244441):
    """Red uppercase section header matching template style."""
    _text(slide, LEFT_MARGIN, top, CONTENT_W, 330200,
          text.upper(), size=13, bold=True, color=ACCENT_RED)


def _subtitle(slide, text, top=533400):
    _text(slide, LEFT_MARGIN, top, CONTENT_W, 461665,
          text, size=12, bold=False, color=GRAY)


def _bullets(slide, left, top, width, height, items, size=9, color=GRAY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        if ": " in item:
            parts = item.split(": ", 1)
            prefix_text = parts[0].replace("**", "")
            rest_text = parts[1]
            run1 = p.add_run()
            run1.text = safe(prefix_text + ": ")
            run1.font.bold = True
            run1.font.size = Pt(size)
            run1.font.color.rgb = DARK_NAVY
            run1.font.name = FONT
            run2 = p.add_run()
            run2.text = safe(rest_text)
            run2.font.size = Pt(size)
            run2.font.color.rgb = color
            run2.font.name = FONT
        else:
            p.text = safe(item)
            p.font.size = Pt(size)
            p.font.color.rgb = color
            p.font.name = FONT
    return txBox


def _card(slide, left, top, width, height, accent_color=None):
    """Card with optional left accent bar."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = CARD_BORDER
    card.line.width = Pt(0.5)
    if accent_color:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, 64008, height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent_color
        bar.line.fill.background()
    return card


def _score_badge(slide, left, top, label, value, max_val=5.0, color=None, with_card=True, card_w=1554480, card_h=457200):
    """Score badge: optional card background + label + value."""
    if color is None:
        if value >= 4.0: color = GREEN
        elif value >= 3.0: color = BLUE
        elif value >= 2.0: color = AMBER
        else: color = RED
    if with_card:
        _card(slide, left, top, card_w, card_h)
    label_left = left + 137160 if with_card else left
    _text(slide, label_left, top + 25400, card_w - 137160 if with_card else 1651000, 152400,
          label, size=10, bold=True, color=DARK_NAVY)
    _text(slide, label_left, top + 165100, 1270000, 152400,
          f"{value:.1f} / {max_val:.0f}", size=14, bold=True, color=color)


def _table(slide, left, top, width, rows_data, size=8, row_height=None):
    n_rows = len(rows_data)
    n_cols = len(rows_data[0]) if rows_data else 1
    row_h = Emu(row_height) if row_height else Emu(228600)
    table_h = row_h * n_rows
    ts = slide.shapes.add_table(n_rows, n_cols, left, top, width, table_h)
    tbl = ts.table
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = safe(str(val))
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(size)
                para.font.name = FONT
                if r == 0:
                    para.font.bold = True
                    para.font.color.rgb = WHITE
                else:
                    para.font.color.rgb = GRAY
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_NAVY
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG
    return ts


# ── Slide builders ────────────────────────────────────────────

def build_title(prs, d, layout):
    """Slide 1: Title on DARK layout. Compact vertical positioning."""
    slide = prs.slides.add_slide(layout)
    # Company name
    _text(slide, LEFT_MARGIN, 1397000, CONTENT_W, 825500,
          safe(d["company_name"]), size=48, bold=True, color=WHITE)
    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, LEFT_MARGIN, 2286000, 1097280, 45720)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_RED
    bar.line.fill.background()
    # Subtitle
    subtitle = d.get("subtitle", "Target Acquisition Memo")
    _text(slide, LEFT_MARGIN, 2438400, 6400800, 731520,
          subtitle, size=18, color=SUBTITLE_C)
    # Quadrant badge
    quad = d.get("quadrant", "")
    qcolor = QUAD_COLORS.get(quad, BLUE)
    qlabel = QUAD_LABELS.get(quad, "")
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    LEFT_MARGIN, 3276600, 2560320, 365760)
    badge.fill.solid()
    badge.fill.fore_color.rgb = qcolor
    badge.line.fill.background()
    tf = badge.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = safe(f"{qlabel}  |  {quad}")
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = FONT
    # Date
    _text(slide, LEFT_MARGIN, 3746500, 4572000, 274320,
          d.get("date", ""), size=9, color=LIGHT_GRAY)
    _add_footer(slide, d["company_name"], d.get("date", ""))


def build_snapshot(prs, d, layout):
    """Slide 2: Company Snapshot & Investment Thesis on LIGHT layout."""
    slide = prs.slides.add_slide(layout)
    _section_header(slide, "Company Snapshot & Investment Thesis", top=274320)
    # KPI table (left) -- starts higher, taller rows
    kpis = d.get("kpis", [])
    if kpis:
        _table(slide, LEFT_MARGIN, 762000, HALF_W, kpis, size=8,
               row_height=int(2971800 / max(len(kpis), 1)))
    # Thesis (right)
    _text(slide, RIGHT_COL, 762000, HALF_W, 274320,
          "INVESTMENT THESIS", size=10, bold=True, color=ACCENT_RED)
    _text(slide, RIGHT_COL, 990600, HALF_W, 2667000,
          d.get("thesis", ""), size=11, color=GRAY)
    # Score badges -- smaller height, tighter spacing
    scores = d.get("scores_summary", {})
    y = 3759200
    x = LEFT_MARGIN
    score_items = list(scores.items())
    for i, (label, val) in enumerate(score_items):
        is_last = (i == len(score_items) - 1)
        cw = 2030526 if is_last else 1554480
        display_label = label
        # Expand abbreviated labels
        if label == "SC_Depth":
            display_label = "Switching Cost Depth"
        _score_badge(slide, x, y, display_label, val, with_card=True,
                     card_w=cw, card_h=457200)
        x += cw + 91440
    # Recommendation
    _text(slide, LEFT_MARGIN, 4216400, CONTENT_W, 152400,
          "RECOMMENDATION", size=10, bold=True, color=ACCENT_RED)
    _text(slide, LEFT_MARGIN, 4381500, CONTENT_W, 431800,
          d.get("recommendation", ""), size=8, bold=False, color=DARK_NAVY)
    _add_footer(slide, d["company_name"], d.get("date", ""))


def build_napkin(prs, d, layout):
    """Slide 3: Vertical Scoring (Napkin) on LIGHT layout."""
    slide = prs.slides.add_slide(layout)
    _section_header(slide, "Vertical Scoring: The Napkin")
    napkin = d.get("napkin", {})
    _subtitle(slide, f"Vertical: {napkin.get('vertical', '')}  |  Composite S_v = {napkin.get('composite', '')}  |  Quadrant: {napkin.get('quadrant', '')}")
    # Table with taller rows
    rows = [["Question", "Score", "Weight", "Rationale"]]
    for q in napkin.get("questions", []):
        rows.append([q["label"], str(q["score"]), str(q["weight"]), q["rationale"]])
    _table(slide, LEFT_MARGIN, 1020465, CONTENT_W, rows, size=8,
           row_height=int(2611120 / max(len(rows), 1)))
    # Evidence -- condensed, positioned after table
    evidence = napkin.get("evidence", [])
    if evidence:
        ev_top = 1020465 + 2611120 + 166381
        _text(slide, LEFT_MARGIN, ev_top, CONTENT_W, 228600,
              "KEY EVIDENCE", size=10, bold=True, color=ACCENT_RED)
        _bullets(slide, LEFT_MARGIN, ev_top + 304798, CONTENT_W, 559127,
                 evidence, size=9, color=GRAY)
    _add_footer(slide, d["company_name"], d.get("date", ""))


def build_company_scores(prs, d, layout):
    """Slide 4: Company-Level Subscores on LIGHT layout. No subtitle."""
    slide = prs.slides.add_slide(layout)
    _section_header(slide, "Company-Level Subscores: Wind / Sail / Fit / AI Signal", top=307340)
    # Table starts higher (no subtitle)
    subs = d.get("company_subscores", [])
    rows = [["Subscore", "Score", "Definition", "Evidence"]]
    for s in subs:
        rows.append([s["name"], str(s["score"]), s["definition"], s["evidence"]])
    _table(slide, LEFT_MARGIN, 762000, CONTENT_W, rows, size=8,
           row_height=int(3114040 / max(len(rows), 1)))
    # Score text (no card backgrounds) -- aligned at bottom
    y_label = 3952240
    y_value = 4130040
    x = LEFT_MARGIN
    for s in subs:
        color = None
        if s["name"] == "Wind": color = GREEN if s["score"] >= 3.5 else AMBER
        elif s["name"] == "Sail": color = AMBER if s["score"] < 3.0 else BLUE
        elif s["name"] == "AI Signal": color = GREEN if s["score"] < 3.0 else AMBER  # low is good
        _text(slide, x, y_label, 1651000, 177800,
              s["name"], size=10, bold=True, color=DARK_NAVY)
        if color is None:
            if s["score"] >= 4.0: color = GREEN
            elif s["score"] >= 3.0: color = BLUE
            elif s["score"] >= 2.0: color = AMBER
            else: color = RED
        _text(slide, x, y_value, 1270000, 177800,
              f"{s['score']:.1f} / 5", size=14, bold=True, color=color)
        x += 1828800
    _add_footer(slide, d["company_name"], d.get("date", ""))


def build_software_moat(prs, d, layout):
    """Slide 5: Software Moat Assessment on LIGHT layout."""
    slide = prs.slides.add_slide(layout)
    _section_header(slide, "Software Moat Assessment", top=127000)
    # Tier as subtitle near top
    moat = d.get("software_moat", {})
    tier = moat.get("tier", "")
    _text(slide, LEFT_MARGIN - 44402, 353367, CONTENT_W, 461665,
          f"Software Moat Tier: {tier}", size=12, bold=False, color=DARK_NAVY)
    # SC_Depth table
    sc = moat.get("sc_depth", [])
    if sc:
        rows = [["Dimension", "Weight", "Score", "Rationale"]]
        for item in sc:
            rows.append([item["dimension"], item["weight"], str(item["score"]), item["rationale"]])
        _table(slide, LEFT_MARGIN, 815032, CONTENT_W, rows, size=8,
               row_height=int(3454400 / max(len(rows), 1)))
    # Composite score badges at bottom with cards
    badges = moat.get("composite_scores", {})
    badge_items = list(badges.items())
    x = LEFT_MARGIN - 44402
    for i, (label, val) in enumerate(badge_items):
        y = 4358232 + (i * 17048)  # slight stagger as in edited version
        _score_badge(slide, x, y, label, val, with_card=True,
                     card_w=1554480, card_h=483474)
        x += 1668120
    _add_footer(slide, d["company_name"], d.get("date", ""))


def build_data_moat(prs, d, layout):
    """Slide 6: Data Moat & Platform Analysis on LIGHT layout."""
    slide = prs.slides.add_slide(layout)
    _section_header(slide, "Data Moat & Platform Analysis", top=274320)
    dm = d.get("data_moat", {})
    content_top = 822960
    # Left column: Data moat
    _text(slide, LEFT_MARGIN, content_top, HALF_W, 274320,
          "DATA MOAT CLASSIFICATION", size=10, bold=True, color=ACCENT_RED)
    _bullets(slide, LEFT_MARGIN, content_top + 320040, HALF_W, 2236510,
             dm.get("classification", []), size=9, color=GRAY)
    # Right column: Platform
    _text(slide, RIGHT_COL, content_top - 5584, HALF_W, 274320,
          "PLATFORM / ECOSYSTEM", size=10, bold=True, color=ACCENT_RED)
    _bullets(slide, RIGHT_COL, content_top + 341888, HALF_W, 1828800,
             dm.get("platform", []), size=9, color=GRAY)
    # AI trajectory bottom
    _text(slide, LEFT_MARGIN, 3383280, CONTENT_W, 274320,
          "AI MOAT TRAJECTORY", size=10, bold=True, color=ACCENT_RED)
    _bullets(slide, LEFT_MARGIN, 3520440, CONTENT_W, 1097280,
             dm.get("ai_trajectory", []), size=9, color=GRAY)
    _add_footer(slide, d["company_name"], d.get("date", ""))


def build_competitive(prs, d, layout):
    """Slide 7: Competitive Landscape on LIGHT layout."""
    slide = prs.slides.add_slide(layout)
    _section_header(slide, "Competitive Landscape", top=274320)
    comp = d.get("competitive", {})
    rows = comp.get("table", [])
    if rows:
        _table(slide, LEFT_MARGIN, 762000, CONTENT_W, rows, size=7,
               row_height=int(1828800 / max(len(rows), 1)))
    notes = comp.get("notes", [])
    if notes:
        table_bottom = 762000 + 1828800
        y = table_bottom + 518160
        _bullets(slide, LEFT_MARGIN, y, CONTENT_W, 1463040, notes, size=9, color=GRAY)
    _add_footer(slide, d["company_name"], d.get("date", ""))


def build_genai(prs, d, layout):
    """Slide 8: GenAI Disruption Assessment on LIGHT layout."""
    slide = prs.slides.add_slide(layout)
    _section_header(slide, "GenAI Disruption Assessment", top=274320)
    genai = d.get("genai_risk", {})
    rows = genai.get("table", [])
    # Append partnership_dependency row if present
    pdep = d.get("partnership_dependency")
    if pdep and rows:
        rows.append(["Partnership Dependency",
                      f"{pdep.get('score', 'N/A')}/5",
                      pdep.get('partner', ''),
                      f"Rev dep: {pdep.get('revenue_dependency', 'N/A')}. {pdep.get('mitigation', '')}"])
    # Append pricing_model_risk row if present
    pmr = d.get("pricing_model_risk")
    if pmr and rows:
        rows.append(["Pricing Model Risk",
                      f"{pmr.get('score', 'N/A')}/5",
                      pmr.get('model', ''),
                      f"Trend: {pmr.get('trend', 'N/A')}. AI impact: {pmr.get('ai_impact', 'N/A')}"])
    if rows:
        _table(slide, LEFT_MARGIN, 698500, CONTENT_W, rows, size=8,
               row_height=int(2494280 / max(len(rows), 1)))
    bullets = genai.get("analysis", [])
    if bullets:
        y = 698500 + 2494280 + 250800 if rows else 762000
        _bullets(slide, LEFT_MARGIN, y, CONTENT_W, 1315745, bullets, size=9, color=GRAY)
    _add_footer(slide, d["company_name"], d.get("date", ""))


def build_100day(prs, d, layout):
    """Slide 9: 100-Day Playbook on LIGHT layout.
    Renders from playbook_100day_structured (phases table) if available,
    falls back to flat playbook_100day array as bullet list."""
    slide = prs.slides.add_slide(layout)
    _section_header(slide, "100-Day Playbook", top=274320)
    structured = d.get("playbook_100day_structured", {})
    phases = structured.get("phases", []) if isinstance(structured, dict) else []
    if phases:
        rows = [["Phase", "Days", "Key Items", "Gate"]]
        for p in phases:
            items_str = "; ".join(p.get("items", [])[:4])  # truncate for space
            if len(p.get("items", [])) > 4:
                items_str += "; ..."
            rows.append([
                f"{p.get('phase', '')}. {p.get('name', '')}",
                p.get("days", ""),
                items_str,
                p.get("gate", "")
            ])
        _table(slide, LEFT_MARGIN, 762000, CONTENT_W, rows, size=7,
               row_height=int(3200000 / max(len(rows), 1)))
    else:
        # Fallback to flat array
        flat = d.get("playbook_100day", [])
        if flat:
            _bullets(slide, LEFT_MARGIN, 762000, CONTENT_W, 3657600,
                     flat, size=9, color=GRAY)
        else:
            _text(slide, LEFT_MARGIN, 762000, CONTENT_W, 457200,
                  "100-Day Playbook to be developed during diligence.",
                  size=11, color=GRAY)
    _add_footer(slide, d["company_name"], d.get("date", ""))


def build_sources(prs, d, layout):
    """Slide 10: Sources & Citations on LIGHT layout."""
    slide = prs.slides.add_slide(layout)
    _section_header(slide, "Sources & Citations", top=274320)
    _bullets(slide, LEFT_MARGIN, 762000, CONTENT_W, 3657600,
             d.get("sources", []), size=7, color=LIGHT_GRAY)
    _add_footer(slide, d["company_name"], d.get("date", ""))


def generate(data_path, output_path, template_path=None):
    with open(data_path) as f:
        d = json.load(f)

    # Use bundled template if none specified and it exists
    if template_path is None and DEFAULT_TEMPLATE.exists():
        template_path = str(DEFAULT_TEMPLATE)

    if template_path:
        prs = Presentation(template_path)
        # Delete existing content slides but keep layouts
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
    else:
        prs = Presentation()
        prs.slide_width = 9144000
        prs.slide_height = 5143500

    # Get layouts (DARK=index 1, LIGHT=index 2, fallback to 0)
    layouts = {l.name: l for l in prs.slide_layouts}
    dark = layouts.get("DARK", prs.slide_layouts[min(1, len(prs.slide_layouts)-1)])
    light = layouts.get("LIGHT", prs.slide_layouts[min(2, len(prs.slide_layouts)-1)])

    # 10-slide structure
    build_title(prs, d, dark)
    build_snapshot(prs, d, light)
    build_napkin(prs, d, light)
    build_company_scores(prs, d, light)
    build_software_moat(prs, d, light)
    build_data_moat(prs, d, light)
    build_competitive(prs, d, light)
    build_genai(prs, d, light)
    build_100day(prs, d, light)
    build_sources(prs, d, light)

    prs.save(output_path)
    print(f"Saved: {output_path} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python3 gen_target_memo_pptx.py data.json output.pptx [template.pptx]")
        sys.exit(1)
    template = sys.argv[3] if len(sys.argv) > 3 else None
    generate(sys.argv[1], sys.argv[2], template)
